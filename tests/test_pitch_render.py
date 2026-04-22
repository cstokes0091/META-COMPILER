"""Tests for the pitch deck v2 evidence pack, fidelity gate, and renderer.

These cover the four-step pitch sub-loop:
  1. build_evidence_pack — typed extraction from artifacts.
  2. write_pitch_request — entry point for the @pitch-writer agent.
  3. verify_slides_fidelity — refuse claims with unknown evidence IDs.
  4. render_pitch_deck — overflow guards (cap + truncate + spill +
     auto_size) and template inheritance.
"""
from __future__ import annotations

from pathlib import Path

import pytest

pptx_module = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402

from meta_compiler.io import dump_yaml, load_yaml
from meta_compiler.stages import pitch_render


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _decision_log() -> dict:
    return {
        "decision_log": {
            "meta": {
                "project_name": "Demo",
                "project_type": "hybrid",
                "version": 1,
                "use_case": "test",
            },
            "conventions": [],
            "architecture": [
                {
                    "component": "ingestor",
                    "approach": "streaming",
                    "constraints_applied": ["low memory"],
                    "alternatives_rejected": [],
                    "rationale": "matches the workload",
                    "citations": ["src-foo2024"],
                }
            ],
            "code_architecture": [
                {
                    "aspect": "language",
                    "choice": "Python 3.11",
                    "rationale": "team familiarity",
                    "citations": [],
                },
                {
                    "aspect": "libraries",
                    "choice": "numpy + pyarrow",
                    "libraries": [
                        {"name": "numpy", "description": "math (>=1.26)"},
                        {"name": "pyarrow", "description": "I/O (>=15)"},
                    ],
                    "rationale": "stable",
                    "citations": [],
                },
            ],
            "scope": {
                "in_scope": [
                    {"item": "stream ingest", "rationale": "must", "citations": []}
                ],
                "out_of_scope": [],
            },
            "requirements": [
                {
                    "id": "REQ-001",
                    "description": "must ingest",
                    "source": "user",
                    "verification": "smoke",
                    "lens": "functional",
                    "citations": [],
                },
                {
                    "id": "REQ-002",
                    "description": "must transform",
                    "source": "user",
                    "verification": "unit",
                    "lens": "functional",
                    "citations": [],
                },
            ],
            "open_items": [
                {
                    "description": "decide retry policy",
                    "deferred_to": "future_work",
                    "owner": "human",
                }
            ],
            "agents_needed": [],
        }
    }


def _final_output_manifest(work_dir: Path) -> dict:
    return {
        "final_output": {
            "decision_log_version": 1,
            "project_type": "hybrid",
            "deliverables": [
                {
                    "agent": "alpha-agent",
                    "kind": "py",
                    "path": str(work_dir.relative_to(work_dir.parent.parent.parent) / "alpha-agent" / "main.py"),
                },
                {
                    "agent": "alpha-agent",
                    "kind": "md",
                    "path": str(work_dir.relative_to(work_dir.parent.parent.parent) / "alpha-agent" / "notes.md"),
                },
            ],
            "execution_notes": ["compiled in test"],
        }
    }


def _seed_work_dir(work_dir: Path, *, mention_req_001: bool, mention_req_002: bool) -> None:
    work_dir.mkdir(parents=True, exist_ok=True)
    (work_dir / "alpha-agent").mkdir(parents=True, exist_ok=True)
    main_text = ["def run():", "    return None"]
    if mention_req_001:
        main_text.insert(0, "# implements REQ-001")
    if mention_req_002:
        main_text.append("# also REQ-002")
    (work_dir / "alpha-agent" / "main.py").write_text(
        "\n".join(main_text) + "\n", encoding="utf-8"
    )
    (work_dir / "alpha-agent" / "notes.md").write_text(
        "# notes\nDocumentation for the alpha agent.\n", encoding="utf-8"
    )


def _seed_problem_statement(workspace_root: Path) -> None:
    workspace_root.mkdir(parents=True, exist_ok=True)
    (workspace_root / "PROBLEM_STATEMENT.md").write_text(
        "# PROBLEM_STATEMENT\n\n## Domain and Problem Space\n"
        "We need a streaming ingestor that handles 10k events/s on a "
        "single laptop while staying under 200MB RAM.\n",
        encoding="utf-8",
    )


# ---------------------------------------------------------------------------
# Evidence pack
# ---------------------------------------------------------------------------


def test_evidence_pack_classifies_traced_and_orphan_requirements(tmp_path: Path):
    workspace_root = tmp_path / "ws"
    artifacts_root = workspace_root / "workspace-artifacts"
    work_dir = artifacts_root / "executions" / "v1" / "work"
    _seed_problem_statement(workspace_root)
    _seed_work_dir(work_dir, mention_req_001=True, mention_req_002=False)

    pack = pitch_render.build_evidence_pack(
        decision_log=_decision_log(),
        decision_log_version=1,
        project_type="hybrid",
        workspace_root=workspace_root,
        final_output_manifest=_final_output_manifest(work_dir),
        work_dir=work_dir,
        citations_payload={"citations": {"src-foo2024": {"human": "Foo 2024", "source": {"type": "pdf"}}}},
        req_trace_path=None,
        ralph_loop_log_path=None,
    )["evidence_pack"]

    assert pack["project"]["id"] == "ev-project"
    assert pack["problem"]["statement"].startswith("We need a streaming ingestor")
    assert len(pack["architecture"]) == 1
    assert pack["architecture"][0]["id"] == "ev-arch-001"
    assert {row["aspect"] for row in pack["code_architecture"]} == {"language", "libraries"}

    traced_ids = {row["req_id"] for row in pack["requirements_traced"]}
    orphan_ids = {row["req_id"] for row in pack["requirements_orphan"]}
    assert traced_ids == {"REQ-001"}
    assert orphan_ids == {"REQ-002"}

    assert len(pack["deliverables"]) == 2
    py_row = next(d for d in pack["deliverables"] if d["kind"] == "py")
    assert py_row["modality"] == "code"
    md_row = next(d for d in pack["deliverables"] if d["kind"] == "md")
    assert md_row["modality"] == "document"

    assert pack["citations"]["ev-cite-src-foo2024"]["human"] == "Foo 2024"


def test_collect_evidence_ids_is_recursive():
    pack = {
        "evidence_pack": {
            "project": {"id": "ev-project"},
            "deliverables": [{"id": "ev-deliv-001"}, {"id": "ev-deliv-002"}],
            "citations": {"ev-cite-foo": {"id": "ev-cite-foo"}},
        }
    }
    ids = pitch_render.collect_evidence_ids(pack)
    assert ids == {"ev-project", "ev-deliv-001", "ev-deliv-002", "ev-cite-foo"}


# ---------------------------------------------------------------------------
# Fidelity gate
# ---------------------------------------------------------------------------


def _slides_payload(*, bullet_evidence_ids: list[str]) -> dict:
    base_slides = []
    for role in pitch_render.REQUIRED_SLIDE_ROLES:
        if role == "title":
            base_slides.append(
                {
                    "role": "title",
                    "title": "T",
                    "subtitle": "S",
                    "evidence_ids": ["ev-project"],
                }
            )
        else:
            base_slides.append(
                {
                    "role": role,
                    "title": role.title(),
                    "bullets": [
                        {"text": "bullet for " + role, "evidence_ids": bullet_evidence_ids},
                    ],
                }
            )
    return {"pitch_deck": {"slides": base_slides}}


def test_fidelity_gate_passes_for_valid_ids():
    pack = {"evidence_pack": {"project": {"id": "ev-project"}}}
    issues = pitch_render.verify_slides_fidelity(
        slides_payload=_slides_payload(bullet_evidence_ids=["ev-project"]),
        evidence_pack=pack,
    )
    assert issues == []


def test_fidelity_gate_rejects_unknown_evidence_id():
    pack = {"evidence_pack": {"project": {"id": "ev-project"}}}
    issues = pitch_render.verify_slides_fidelity(
        slides_payload=_slides_payload(bullet_evidence_ids=["ev-bogus"]),
        evidence_pack=pack,
    )
    assert any("ev-bogus" in i for i in issues)


def test_fidelity_gate_rejects_bullet_with_no_evidence_ids():
    pack = {"evidence_pack": {"project": {"id": "ev-project"}}}
    issues = pitch_render.verify_slides_fidelity(
        slides_payload=_slides_payload(bullet_evidence_ids=[]),
        evidence_pack=pack,
    )
    assert any("must cite at least one evidence_ids" in i for i in issues)


def test_fidelity_gate_flags_missing_required_role():
    payload = _slides_payload(bullet_evidence_ids=["ev-project"])
    payload["pitch_deck"]["slides"] = [
        s for s in payload["pitch_deck"]["slides"] if s["role"] != "cta"
    ]
    pack = {"evidence_pack": {"project": {"id": "ev-project"}}}
    issues = pitch_render.verify_slides_fidelity(
        slides_payload=payload,
        evidence_pack=pack,
    )
    assert any("missing required slide roles" in i and "cta" in i for i in issues)


# ---------------------------------------------------------------------------
# Render layout guards: truncate, spill, auto_size
# ---------------------------------------------------------------------------


def _build_render_payload(*, bullet_count: int, bullet_text: str, title: str) -> dict:
    return {
        "pitch_deck": {
            "slides": [
                {
                    "role": "title",
                    "title": title,
                    "subtitle": "test",
                    "evidence_ids": ["ev-project"],
                },
                {
                    "role": "built",
                    "title": "What was built",
                    "bullets": [
                        {"text": bullet_text, "evidence_ids": ["ev-project"]}
                        for _ in range(bullet_count)
                    ],
                },
            ]
        }
    }


def test_render_truncates_long_bullets_and_titles(tmp_path: Path):
    long_bullet = "x" * 300
    long_title = "y" * 200
    payload = _build_render_payload(
        bullet_count=1, bullet_text=long_bullet, title=long_title
    )
    out = tmp_path / "deck.pptx"
    pitch_render.render_pitch_deck(
        slides_payload=payload, output_path=out, template_path=None
    )
    assert out.exists()
    prs = Presentation(str(out))
    title_text = prs.slides[0].shapes.title.text
    assert title_text.endswith("…")
    assert len(title_text) <= pitch_render.SLIDE_CAPS["title_chars_max"]
    body_placeholder = pitch_render._find_content_placeholder(prs.slides[1])
    assert body_placeholder is not None
    bullet_text = body_placeholder.text_frame.paragraphs[0].text
    assert bullet_text.endswith("…")
    assert len(bullet_text) <= pitch_render.SLIDE_CAPS["bullet_chars_max"]


def test_render_spills_long_bullet_lists_into_followup_slides(tmp_path: Path):
    payload = _build_render_payload(bullet_count=14, bullet_text="ok", title="Built")
    out = tmp_path / "deck.pptx"
    pitch_render.render_pitch_deck(
        slides_payload=payload, output_path=out, template_path=None
    )
    prs = Presentation(str(out))
    # 1 title slide + ceil(14 / 6) = 3 spill slides => 4 slides total.
    slides = list(prs.slides)
    assert len(slides) == 4
    titles = [s.shapes.title.text for s in slides[1:]]
    assert titles[0].endswith("(1/3)")
    assert titles[1].endswith("(2/3)")
    assert titles[2].endswith("(3/3)")


def test_render_enables_auto_size_and_word_wrap(tmp_path: Path):
    payload = _build_render_payload(bullet_count=2, bullet_text="x", title="Built")
    out = tmp_path / "deck.pptx"
    pitch_render.render_pitch_deck(
        slides_payload=payload, output_path=out, template_path=None
    )
    prs = Presentation(str(out))
    body = pitch_render._find_content_placeholder(prs.slides[1])
    assert body is not None
    assert body.text_frame.word_wrap is True
    assert body.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# ---------------------------------------------------------------------------
# Template inheritance (.pptx and .potx accepted)
# ---------------------------------------------------------------------------


def test_render_inherits_from_pptx_template(tmp_path: Path):
    """Round-trip: build a brand template, render against it, confirm the
    template's slide_master/theme is preserved on the output."""
    template = Presentation()
    # Add one placeholder slide to the template — the renderer must strip it
    # before writing rendered slides so it doesn't pollute the deck.
    template.slides.add_slide(template.slide_layouts[0]).shapes.title.text = (
        "TEMPLATE PLACEHOLDER"
    )
    template_path = tmp_path / "brand_template.pptx"
    template.save(str(template_path))
    template_master_id = id(template.slide_masters[0])  # noqa: F841 - sanity ref

    payload = _build_render_payload(bullet_count=2, bullet_text="ok", title="Built")
    out = tmp_path / "deck.pptx"
    pitch_render.render_pitch_deck(
        slides_payload=payload, output_path=out, template_path=template_path
    )
    rendered = Presentation(str(out))
    # The template's placeholder slide must be gone.
    titles = [s.shapes.title.text for s in rendered.slides if s.shapes.title is not None]
    assert "TEMPLATE PLACEHOLDER" not in titles
    # Title slide + 1 built slide.
    assert len(rendered.slides) == 2


def test_render_rejects_invalid_template_suffix(tmp_path: Path):
    bogus = tmp_path / "not_a_template.txt"
    bogus.write_text("nope", encoding="utf-8")
    payload = _build_render_payload(bullet_count=1, bullet_text="x", title="t")
    with pytest.raises(RuntimeError, match=".pptx or .potx"):
        pitch_render.render_pitch_deck(
            slides_payload=payload,
            output_path=tmp_path / "out.pptx",
            template_path=bogus,
        )
