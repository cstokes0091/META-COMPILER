#!/usr/bin/env python3
"""Extract text from common document formats (docx, xlsx, pptx, pdf).

Usage:
    python scripts/read_document.py <file_path> [--output <output_path>]

Reads the input file and prints extracted text to stdout.  When --output is
given the text is written to that path instead.

Supported formats: .docx, .xlsx, .pptx, .pdf, .txt, .md
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def read_docx(path: Path) -> str:
    """Extract full text from a .docx file."""
    from docx import Document

    doc = Document(str(path))
    lines: list[str] = []
    for para in doc.paragraphs:
        lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def read_xlsx(path: Path) -> str:
    """Extract full text from a .xlsx file."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append(" | ".join(cells))
        lines.append("")
    wb.close()
    return "\n".join(lines)


def read_pptx(path: Path) -> str:
    """Extract full text from a .pptx file."""
    from pptx import Presentation

    prs = Presentation(str(path))
    lines: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
        lines.append("")
    return "\n".join(lines)


def read_pdf(path: Path) -> str:
    """Extract full text from a .pdf file."""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    lines: list[str] = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        lines.append(f"## Page {page_num + 1}")
        lines.append(text.strip())
        lines.append("")
    doc.close()
    return "\n".join(lines)


def read_plain(path: Path) -> str:
    """Read a plain text or markdown file."""
    return path.read_text(encoding="utf-8", errors="replace")


READERS = {
    ".docx": read_docx,
    ".xlsx": read_xlsx,
    ".pptx": read_pptx,
    ".pdf": read_pdf,
    ".txt": read_plain,
    ".md": read_plain,
    ".rst": read_plain,
    ".tex": read_plain,
    ".csv": read_plain,
}


def read_document(path: Path) -> str:
    """Read a document file and return its full text."""
    suffix = path.suffix.lower()
    reader = READERS.get(suffix)
    if reader is None:
        raise ValueError(
            f"Unsupported file format: {suffix}. "
            f"Supported formats: {', '.join(sorted(READERS))}"
        )
    return reader(path)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Extract text from documents (docx, xlsx, pptx, pdf, txt, md)."
    )
    parser.add_argument("file", help="Path to the document file")
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="Write extracted text to this file instead of stdout",
    )
    args = parser.parse_args(argv)

    file_path = Path(args.file)
    if not file_path.exists():
        print(f"ERROR: File not found: {file_path}", file=sys.stderr)
        return 1

    try:
        text = read_document(file_path)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1

    if args.output:
        Path(args.output).write_text(text, encoding="utf-8")
        print(f"Extracted text written to {args.output}")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
