"""Pitch deck v2 — typed evidence, LLM-drafted, template-aware.

The pitch sub-loop runs in four steps controlled by `phase4-finalize
--pitch-step`:

  1. evidence  — `build_evidence_pack()` extracts a typed pack of project
                 facts (problem, architecture, deliverables, REQ trace, …)
                 from the Decision Log + FINAL_OUTPUT_MANIFEST + work dir.
                 Writes `runtime/phase4/evidence_pack.yaml` and the
                 `pitch_request.yaml` the pitch-writer agent reads.
  2. draft     — handed to the `@pitch-writer` agent in Copilot. The agent
                 reads the evidence pack and writes
                 `runtime/phase4/slides.yaml` whose every bullet cites at
                 least one `evidence_ids[...]`.
  3. verify    — `verify_slides_fidelity()` walks every bullet's
                 `evidence_ids` and refuses to render unless every ID
                 resolves to a known evidence entry.
  4. render    — `render_pitch_deck()` walks `slides.yaml` and renders to
                 `pitches/pitch_v{N}.pptx` using strict layout guards
                 (cap + auto-shrink + spill) and an optional `.pptx` /
                 `.potx` template.
"""
from __future__ import annotations

import re
from collections import defaultdict
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore
except ImportError:  # pragma: no cover - exercised when dependency missing
    Presentation = None  # type: ignore[assignment]
    MSO_AUTO_SIZE = None  # type: ignore[assignment]

from ..io import dump_yaml, load_yaml
from ..utils import iso_now, read_text_safe, slugify


# ---------------------------------------------------------------------------
# Layout guards — every constant is intentionally conservative; the renderer
# truncates and spills aggressively rather than letting any slide overflow.
# ---------------------------------------------------------------------------

SLIDE_CAPS: dict[str, int] = {
    "title_chars_max": 70,
    "subtitle_chars_max": 200,
    "bullet_chars_max": 140,
    "bullets_per_slide_max": 6,
    "speaker_notes_chars_max": 600,
}

# Project-pitch arc — the canonical 7-slide structure. The pitch-writer agent
# is expected to emit exactly these roles in this order. Verify rejects a
# slides.yaml that omits any required role.
REQUIRED_SLIDE_ROLES: tuple[str, ...] = (
    "title",
    "problem",
    "approach",
    "built",
    "evidence",
    "why",
    "cta",
)


# ---------------------------------------------------------------------------
# Step 1 — build_evidence_pack
# ---------------------------------------------------------------------------


def build_evidence_pack(
    *,
    decision_log: dict[str, Any],
    decision_log_version: int,
    project_type: str,
    workspace_root: Path,
    final_output_manifest: dict[str, Any],
    work_dir: Path,
    citations_payload: dict[str, Any] | None,
    req_trace_path: Path | None,
    ralph_loop_log_path: Path | None,
    final_dir: Path | None = None,
) -> dict[str, Any]:
    """Extract a typed evidence pack from the Decision Log + execution outputs.

    Every fact gets a stable `id` (e.g. `ev-deliv-003`, `ev-req-007`) so the
    pitch-writer agent can cite specific evidence and the fidelity gate can
    refuse claims whose IDs don't resolve.

    When `final_dir` is provided AND the final-synthesis sub-stage has run
    (`final_output_manifest.synthesis_status == "synthesized"`), the pack
    additionally exposes `assembled_deliverables[]` — high-level
    `ev-final-NNN` evidence IDs that the pitch-writer should prefer on the
    `built` slide. The legacy `deliverables[]` list (`ev-deliv-NNN`) keeps
    pointing at per-capability fragments for the `evidence` slide where
    coverage breadth matters.
    """
    root = decision_log.get("decision_log") or {}
    meta = root.get("meta") or {}

    evidence_pack: dict[str, Any] = {
        "generated_at": iso_now(),
        "decision_log_version": decision_log_version,
        "project_type": project_type,
        "project": _evidence_project(meta),
        "problem": _evidence_problem(workspace_root),
        "architecture": _evidence_architecture(root),
        "code_architecture": _evidence_code_architecture(root),
        "scope": _evidence_scope(root),
        "requirements_traced": [],
        "requirements_orphan": [],
        "deliverables": _evidence_deliverables(
            final_output_manifest=final_output_manifest, work_dir=work_dir
        ),
        "open_items": _evidence_open_items(root),
        "citations": _evidence_citations(citations_payload, root),
        "execution": _evidence_execution(
            final_output_manifest=final_output_manifest,
            ralph_loop_log_path=ralph_loop_log_path,
        ),
    }

    if final_dir is not None and final_dir.exists() and any(final_dir.rglob("*")):
        evidence_pack["assembled_deliverables"] = _evidence_assembled_deliverables(
            final_dir=final_dir,
            workspace_root=workspace_root,
        )

    pack: dict[str, Any] = {"evidence_pack": evidence_pack}

    traced, orphan = _evidence_requirements(
        root=root,
        work_dir=work_dir,
        req_trace_path=req_trace_path,
        final_dir=final_dir,
    )
    pack["evidence_pack"]["requirements_traced"] = traced
    pack["evidence_pack"]["requirements_orphan"] = orphan

    return pack


def _evidence_project(meta: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": "ev-project",
        "name": str(meta.get("project_name") or "META-COMPILER Project"),
        "type": str(meta.get("project_type") or "algorithm"),
        "version": meta.get("version"),
        "use_case": meta.get("use_case") or "",
    }


_PROBLEM_HEADING_RE = re.compile(r"^##+\s*Domain and Problem Space\s*$", re.IGNORECASE)


def _evidence_problem(workspace_root: Path) -> dict[str, Any]:
    statement_path = workspace_root / "PROBLEM_STATEMENT.md"
    if not statement_path.exists():
        return {
            "id": "ev-problem",
            "statement": "",
            "source_path": str(statement_path),
            "missing": True,
        }
    text = read_text_safe(statement_path)
    statement_chunk = _extract_section(text, _PROBLEM_HEADING_RE) or _first_paragraph(text)
    return {
        "id": "ev-problem",
        "statement": statement_chunk.strip(),
        "source_path": str(statement_path),
        "missing": False,
    }


def _extract_section(text: str, heading_re: re.Pattern[str]) -> str:
    """Return the prose under the first heading matching `heading_re`,
    stopping at the next `## ` heading or end-of-file."""
    lines = text.splitlines()
    collecting = False
    body: list[str] = []
    for line in lines:
        if heading_re.match(line):
            collecting = True
            continue
        if collecting and line.startswith("## ") and "Domain and Problem Space" not in line:
            break
        if collecting:
            body.append(line)
    return "\n".join(body).strip()


def _first_paragraph(text: str) -> str:
    blocks = [b.strip() for b in text.split("\n\n") if b.strip() and not b.lstrip().startswith("#")]
    return blocks[0] if blocks else ""


def _evidence_architecture(root: dict[str, Any]) -> list[dict[str, Any]]:
    rows = root.get("architecture") or []
    out: list[dict[str, Any]] = []
    for idx, row in enumerate(rows, start=1):
        if not isinstance(row, dict):
            continue
        out.append(
            {
                "id": f"ev-arch-{idx:03d}",
                "component": row.get("component", ""),
                "approach": row.get("approach", ""),
                "constraints_applied": list(row.get("constraints_applied") or []),
                "alternatives_rejected": [
                    dict(alt) for alt in (row.get("alternatives_rejected") or [])
                ],
                "rationale": row.get("rationale", ""),
                "citations": list(row.get("citations") or []),
            }
        )
    return out


def _evidence_code_architecture(root: dict[str, Any]) -> list[dict[str, Any]]:
    rows = root.get("code_architecture") or []
    out: list[dict[str, Any]] = []
    for idx, row in enumerate(rows, start=1):
        if not isinstance(row, dict):
            continue
        out.append(
            {
                "id": f"ev-codearch-{idx:03d}",
                "aspect": row.get("aspect", ""),
                "choice": row.get("choice", ""),
                "libraries": [dict(lib) for lib in (row.get("libraries") or [])],
                "module_layout": row.get("module_layout"),
                "rationale": row.get("rationale", ""),
                "citations": list(row.get("citations") or []),
            }
        )
    return out


def _evidence_scope(root: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    scope = root.get("scope") or {}
    in_rows = []
    for idx, row in enumerate(scope.get("in_scope") or [], start=1):
        if not isinstance(row, dict):
            continue
        in_rows.append(
            {
                "id": f"ev-scope-in-{idx:03d}",
                "item": row.get("item", ""),
                "rationale": row.get("rationale", ""),
                "citations": list(row.get("citations") or []),
            }
        )
    out_rows = []
    for idx, row in enumerate(scope.get("out_of_scope") or [], start=1):
        if not isinstance(row, dict):
            continue
        out_rows.append(
            {
                "id": f"ev-scope-out-{idx:03d}",
                "item": row.get("item", ""),
                "rationale": row.get("rationale", ""),
                "revisit_if": row.get("revisit_if", ""),
                "citations": list(row.get("citations") or []),
            }
        )
    return {"in_scope": in_rows, "out_of_scope": out_rows}


# Lightweight extension → modality mapping. Code suffixes mark deliverables
# as `code`; everything else is `document`.
_CODE_SUFFIXES = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".rs", ".go", ".java", ".kt",
    ".c", ".cc", ".cpp", ".h", ".hpp", ".rb", ".php", ".cs", ".swift",
    ".sh", ".bash", ".zsh", ".sql", ".m", ".mm", ".scala", ".jl",
}


def _evidence_deliverables(
    *, final_output_manifest: dict[str, Any], work_dir: Path
) -> list[dict[str, Any]]:
    """Per-capability fragments (the `ev-deliv-*` evidence list).

    When the final-synthesis sub-stage has run, the manifest's top-level
    `deliverables[]` is the *assembled* artifact set; the per-capability
    fragments live under `fragments[]`. We prefer that when available so
    `ev-deliv-*` IDs continue to mean "per-capability fragment" — letting
    the pitch-writer use them on the `evidence` slide where coverage
    breadth matters. Pre-synthesis manifests still fall through to
    `deliverables[]` (the legacy shape).
    """
    final_output = final_output_manifest.get("final_output") or {}
    fragments = final_output.get("fragments") or []
    if not fragments:
        fragments = final_output.get("deliverables") or []

    out: list[dict[str, Any]] = []
    for idx, row in enumerate(fragments, start=1):
        if not isinstance(row, dict):
            continue
        path_str = str(row.get("path") or "")
        agent = str(row.get("agent") or "")
        path_obj = Path(path_str)
        absolute = path_obj if path_obj.is_absolute() else (work_dir.parent.parent / path_str)
        size_bytes = 0
        line_count = 0
        try:
            if absolute.exists() and absolute.is_file():
                size_bytes = absolute.stat().st_size
                if absolute.suffix.lower() in _CODE_SUFFIXES or absolute.suffix.lower() in {
                    ".md", ".txt", ".yaml", ".yml", ".json", ".rst",
                }:
                    try:
                        line_count = sum(
                            1 for _ in absolute.read_text(encoding="utf-8").splitlines()
                        )
                    except UnicodeDecodeError:
                        line_count = 0
        except OSError:
            pass

        modality = "code" if absolute.suffix.lower() in _CODE_SUFFIXES else "document"

        out.append(
            {
                "id": f"ev-deliv-{idx:03d}",
                "agent": agent,
                "kind": str(row.get("kind") or "file"),
                "path": path_str,
                "modality": modality,
                "size_bytes": size_bytes,
                "line_count": line_count,
            }
        )
    return out


def _evidence_assembled_deliverables(
    *, final_dir: Path, workspace_root: Path
) -> list[dict[str, Any]]:
    """High-level assembled artifact list (the `ev-final-*` evidence list).

    Walked directly from `executions/v{N}/final/<bucket>/`. The pitch-writer
    cites these IDs on the `built` slide because they're the truthful
    "what we shipped" answer — a single library / document / application —
    rather than the dozens of per-capability fragments under work/.
    """
    out: list[dict[str, Any]] = []
    if not final_dir.exists():
        return out
    for idx, path in enumerate(sorted(final_dir.rglob("*")), start=1):
        if not path.is_file():
            continue
        try:
            relative = path.relative_to(final_dir)
        except ValueError:
            continue
        bucket = relative.parts[0] if relative.parts else ""
        try:
            relative_to_workspace = path.relative_to(workspace_root)
            display_path = str(relative_to_workspace)
        except ValueError:
            display_path = str(path)
        size_bytes = 0
        line_count = 0
        try:
            size_bytes = path.stat().st_size
            if path.suffix.lower() in _CODE_SUFFIXES or path.suffix.lower() in {
                ".md", ".txt", ".yaml", ".yml", ".json", ".rst", ".toml",
            }:
                try:
                    line_count = sum(
                        1 for _ in path.read_text(encoding="utf-8").splitlines()
                    )
                except UnicodeDecodeError:
                    line_count = 0
        except OSError:
            pass
        modality = "code" if path.suffix.lower() in _CODE_SUFFIXES else "document"
        out.append(
            {
                "id": f"ev-final-{idx:03d}",
                "bucket": bucket,
                "kind": path.suffix.lstrip(".") or "file",
                "path": display_path,
                "modality": modality,
                "size_bytes": size_bytes,
                "line_count": line_count,
            }
        )
    return out


def _evidence_open_items(root: dict[str, Any]) -> list[dict[str, Any]]:
    rows = root.get("open_items") or []
    out: list[dict[str, Any]] = []
    for idx, row in enumerate(rows, start=1):
        if not isinstance(row, dict):
            continue
        out.append(
            {
                "id": f"ev-open-{idx:03d}",
                "description": row.get("description", ""),
                "deferred_to": row.get("deferred_to", ""),
                "owner": row.get("owner", ""),
            }
        )
    return out


def _evidence_citations(
    citations_payload: dict[str, Any] | None,
    root: dict[str, Any],
) -> dict[str, dict[str, Any]]:
    """Return a mapping of `ev-cite-<sluggified citation id>` -> citation
    metadata, restricted to citations actually referenced by the Decision Log.
    """
    payload = citations_payload or {}
    citations_index = payload.get("citations") or {}
    referenced: set[str] = set()
    for section_key in ("conventions", "architecture", "code_architecture", "requirements"):
        for row in root.get(section_key) or []:
            if not isinstance(row, dict):
                continue
            for cid in row.get("citations") or []:
                if isinstance(cid, str) and cid.strip():
                    referenced.add(cid.strip())

    out: dict[str, dict[str, Any]] = {}
    for cid in sorted(referenced):
        ev_id = f"ev-cite-{slugify(cid)}"
        entry = citations_index.get(cid) if isinstance(citations_index, dict) else None
        if not isinstance(entry, dict):
            entry = {}
        source_meta = entry.get("source") if isinstance(entry.get("source"), dict) else {}
        out[ev_id] = {
            "id": ev_id,
            "citation_id": cid,
            "human": str(entry.get("human") or cid),
            "source_type": str(source_meta.get("type") or "unknown"),
        }
    return out


def _evidence_execution(
    *,
    final_output_manifest: dict[str, Any],
    ralph_loop_log_path: Path | None,
) -> dict[str, Any]:
    final_output = final_output_manifest.get("final_output") or {}
    notes = list(final_output.get("execution_notes") or [])

    cycle_summary: dict[str, Any] = {
        "total_cycles": 0,
        "force_advanced": [],
        "revise_count": 0,
    }
    if ralph_loop_log_path is not None and ralph_loop_log_path.exists():
        ralph_payload = load_yaml(ralph_loop_log_path) or {}
        ralph_log = ralph_payload.get("ralph_loop_log") if isinstance(ralph_payload, dict) else None
        if isinstance(ralph_log, dict):
            cycle_summary["total_cycles"] = int(ralph_log.get("total_cycles", 0) or 0)
            forced = ralph_log.get("force_advanced") or []
            cycle_summary["force_advanced"] = [str(s) for s in forced if isinstance(s, str)]
            cycle_summary["revise_count"] = int(ralph_log.get("revise_count", 0) or 0)

    return {
        "id": "ev-exec",
        "execution_notes": notes,
        "cycle_summary": cycle_summary,
    }


_REQ_RE = re.compile(r"REQ-\d{3}")


def _evidence_requirements(
    *,
    root: dict[str, Any],
    work_dir: Path,
    req_trace_path: Path | None,
    final_dir: Path | None = None,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Cross-reference every REQ-NNN against the assembled tree first, the
    implementer work files second.

    A requirement is `traced` when at least one assembled-tree file (or, as
    fallback, a work_dir file) mentions its REQ ID. Otherwise it's `orphan`
    — a real accuracy red flag the deck must surface honestly rather than
    paper over.

    Preferring `final_dir/` when present aligns the pitch deck's "this REQ
    was delivered" claim with the *assembled* artifacts the user actually
    sees, not just per-capability fragments under work/.
    """
    requirements = root.get("requirements") or []

    final_files: list[Path] = []
    if final_dir is not None and final_dir.exists():
        for path in final_dir.rglob("*"):
            if path.is_file():
                final_files.append(path)

    work_files: list[Path] = []
    if work_dir.exists():
        for path in work_dir.rglob("*"):
            if path.is_file():
                work_files.append(path)

    file_text_cache: dict[Path, str] = {}

    def _file_mentions(req_id: str, path: Path) -> bool:
        if path not in file_text_cache:
            try:
                file_text_cache[path] = path.read_text(encoding="utf-8")
            except (OSError, UnicodeDecodeError):
                file_text_cache[path] = ""
        return req_id in file_text_cache[path]

    traced: list[dict[str, Any]] = []
    orphan: list[dict[str, Any]] = []
    counter_t = 0
    counter_o = 0

    for req in requirements:
        if not isinstance(req, dict):
            continue
        req_id = str(req.get("id") or "")
        if not _REQ_RE.fullmatch(req_id):
            continue

        # Prefer assembled tree; fall back to work_dir.
        implementing_paths = [p for p in final_files if _file_mentions(req_id, p)]
        scope = "final"
        if not implementing_paths:
            implementing_paths = [p for p in work_files if _file_mentions(req_id, p)]
            scope = "work" if implementing_paths else scope

        def _display(p: Path) -> str:
            for base in (final_dir, work_dir.parent.parent if work_dir else None):
                if base is None:
                    continue
                try:
                    return str(p.relative_to(base))
                except ValueError:
                    continue
            return str(p)

        implementing = [_display(p) for p in implementing_paths]

        if implementing:
            counter_t += 1
            traced.append(
                {
                    "id": f"ev-req-{counter_t:03d}",
                    "req_id": req_id,
                    "description": str(req.get("description", "")),
                    "verification": str(req.get("verification", "")),
                    "lens": str(req.get("lens", "")),
                    "implementing_files": implementing[:8],
                    "trace_scope": scope,
                }
            )
        else:
            counter_o += 1
            orphan.append(
                {
                    "id": f"ev-orphan-{counter_o:03d}",
                    "req_id": req_id,
                    "description": str(req.get("description", "")),
                    "reason": "no work or final file references this REQ ID",
                }
            )

    return traced, orphan


# ---------------------------------------------------------------------------
# Pitch request — the input handed to the @pitch-writer agent
# ---------------------------------------------------------------------------


def write_pitch_request(
    *,
    pitch_request_path: Path,
    evidence_pack_path: Path,
    slides_path: Path,
    pptx_output_path: Path,
    template_path: Path | None,
    decision_log_version: int,
) -> Path:
    payload = {
        "pitch_request": {
            "generated_at": iso_now(),
            "decision_log_version": decision_log_version,
            "evidence_pack_path": str(evidence_pack_path),
            "slides_output_path": str(slides_path),
            "pptx_output_path": str(pptx_output_path),
            "template_path": str(template_path) if template_path else None,
            "slide_caps": SLIDE_CAPS,
            "required_slide_roles": list(REQUIRED_SLIDE_ROLES),
            "instructions": [
                "Read the evidence pack at evidence_pack_path.",
                "Draft a 7-slide project-pitch deck (roles: " + ", ".join(REQUIRED_SLIDE_ROLES) + ").",
                "Every bullet MUST cite at least one evidence_ids[...] from the pack.",
                "Describe THIS project — never the META-COMPILER framework.",
                "Surface requirements_orphan[] honestly in the 'built' or 'why' slide if any are present.",
                "Respect SLIDE_CAPS (the renderer truncates and spills, but anchor your draft to them).",
                "Write the result to slides_output_path (YAML schema: pitch_deck.slides[]).",
                "Then re-run `meta-compiler phase4-finalize --pitch-step=render` to produce the .pptx.",
            ],
        }
    }
    pitch_request_path.parent.mkdir(parents=True, exist_ok=True)
    dump_yaml(pitch_request_path, payload)
    return pitch_request_path


# ---------------------------------------------------------------------------
# Step 3 — verify_slides_fidelity
# ---------------------------------------------------------------------------


def collect_evidence_ids(evidence_pack: dict[str, Any]) -> set[str]:
    """Return the flat set of every `id` string addressable in the pack."""
    pack = evidence_pack.get("evidence_pack") if isinstance(evidence_pack, dict) else None
    if not isinstance(pack, dict):
        return set()

    ids: set[str] = set()

    def _walk(value: Any) -> None:
        if isinstance(value, dict):
            ev_id = value.get("id")
            if isinstance(ev_id, str) and ev_id.startswith("ev-"):
                ids.add(ev_id)
            for v in value.values():
                _walk(v)
        elif isinstance(value, list):
            for item in value:
                _walk(item)

    _walk(pack)
    return ids


def verify_slides_fidelity(
    *,
    slides_payload: dict[str, Any],
    evidence_pack: dict[str, Any],
) -> list[str]:
    """Return a list of violations. Empty list = render is allowed.

    Two kinds of failures:
      - structural: missing required slide roles, malformed slides.
      - fidelity: a bullet's `evidence_ids[...]` references an unknown ID.
    """
    issues: list[str] = []
    deck = slides_payload.get("pitch_deck") if isinstance(slides_payload, dict) else None
    if not isinstance(deck, dict):
        issues.append("pitch_deck: missing root object")
        return issues

    slides = deck.get("slides") or []
    if not isinstance(slides, list) or not slides:
        issues.append("pitch_deck.slides: must be a non-empty list")
        return issues

    valid_ids = collect_evidence_ids(evidence_pack)

    seen_roles: set[str] = set()
    for idx, slide in enumerate(slides):
        if not isinstance(slide, dict):
            issues.append(f"pitch_deck.slides[{idx}]: must be an object")
            continue
        role = slide.get("role")
        if not isinstance(role, str) or not role:
            issues.append(f"pitch_deck.slides[{idx}]: missing 'role'")
        else:
            seen_roles.add(role)
        title = slide.get("title")
        if not isinstance(title, str) or not title.strip():
            issues.append(f"pitch_deck.slides[{idx}] (role={role}): missing 'title'")

        evidence_ids = slide.get("evidence_ids") or []
        for ev_id in evidence_ids:
            if not isinstance(ev_id, str) or not ev_id.startswith("ev-"):
                issues.append(
                    f"pitch_deck.slides[{idx}].evidence_ids: '{ev_id}' is not a valid evidence id"
                )
                continue
            if ev_id not in valid_ids:
                issues.append(
                    f"pitch_deck.slides[{idx}].evidence_ids: '{ev_id}' does not resolve to "
                    "any entry in evidence_pack"
                )

        bullets = slide.get("bullets") or []
        if role in {"title"}:
            # title slide uses subtitle, not bullets
            continue
        if not isinstance(bullets, list):
            issues.append(f"pitch_deck.slides[{idx}]: 'bullets' must be a list")
            continue
        for jdx, bullet in enumerate(bullets):
            if not isinstance(bullet, dict):
                issues.append(
                    f"pitch_deck.slides[{idx}].bullets[{jdx}]: must be {{text, evidence_ids[]}}"
                )
                continue
            text = bullet.get("text")
            if not isinstance(text, str) or not text.strip():
                issues.append(
                    f"pitch_deck.slides[{idx}].bullets[{jdx}]: missing 'text'"
                )
            bev_ids = bullet.get("evidence_ids") or []
            if not isinstance(bev_ids, list) or not bev_ids:
                issues.append(
                    f"pitch_deck.slides[{idx}].bullets[{jdx}]: must cite at least one "
                    "evidence_ids[...] entry"
                )
                continue
            for bev_id in bev_ids:
                if not isinstance(bev_id, str) or not bev_id.startswith("ev-"):
                    issues.append(
                        f"pitch_deck.slides[{idx}].bullets[{jdx}].evidence_ids: "
                        f"'{bev_id}' is not a valid evidence id"
                    )
                    continue
                if bev_id not in valid_ids:
                    issues.append(
                        f"pitch_deck.slides[{idx}].bullets[{jdx}].evidence_ids: "
                        f"'{bev_id}' does not resolve to any entry in evidence_pack"
                    )

    missing_roles = [r for r in REQUIRED_SLIDE_ROLES if r not in seen_roles]
    if missing_roles:
        issues.append(
            f"pitch_deck: missing required slide roles {missing_roles}"
        )

    return issues


# ---------------------------------------------------------------------------
# Step 4 — render_pitch_deck
# ---------------------------------------------------------------------------


def _truncate(text: str, max_chars: int) -> str:
    text = text or ""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _spill_bullets(
    bullets: list[dict[str, Any]],
    bullets_per_slide_max: int,
) -> list[list[dict[str, Any]]]:
    if not bullets:
        return [[]]
    out: list[list[dict[str, Any]]] = []
    for start in range(0, len(bullets), bullets_per_slide_max):
        out.append(bullets[start : start + bullets_per_slide_max])
    return out


def _resolve_layout(prs, layout_name: str | None, fallback_index: int):
    if layout_name:
        for layout in prs.slide_layouts:
            if str(layout.name).strip().lower() == layout_name.strip().lower():
                return layout
    if fallback_index < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def _remove_existing_slides(prs) -> int:
    """Strip slides shipped with the template so the rendered deck is clean.

    `python-pptx` doesn't expose a direct API; we drop slide IDs from the
    presentation's sldIdLst and unlink the associated parts.
    """
    removed = 0
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001 - intentional internal access
    rId_to_drop: list[str] = []
    for sldId in list(sldIdLst):
        rId_to_drop.append(sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ))
        sldIdLst.remove(sldId)
        removed += 1
    for rId in rId_to_drop:
        if rId is None:
            continue
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    return removed


def _set_text_frame(
    text_frame,
    *,
    text: str,
    max_chars: int,
) -> None:
    """Set a text frame's body to a single paragraph with overflow guards."""
    text_frame.clear()
    if MSO_AUTO_SIZE is not None:
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except (AttributeError, ValueError):  # pragma: no cover - defensive
            pass
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.text = _truncate(text, max_chars)


def _set_bulleted_text_frame(
    text_frame,
    *,
    bullets: list[str],
    bullet_chars_max: int,
) -> None:
    text_frame.clear()
    if MSO_AUTO_SIZE is not None:
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except (AttributeError, ValueError):  # pragma: no cover - defensive
            pass
    text_frame.word_wrap = True
    if not bullets:
        text_frame.paragraphs[0].text = ""
        return
    for idx, raw in enumerate(bullets):
        text = _truncate(raw, bullet_chars_max)
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
            paragraph.text = text
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.text = text
        paragraph.level = 0


def render_pitch_deck(
    *,
    slides_payload: dict[str, Any],
    output_path: Path,
    template_path: Path | None,
) -> Path:
    """Render slides.yaml -> pptx with strict layout guards."""
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required for Stage 4 pitch rendering. "
            "Install with: python3 -m pip install -r requirements.txt"
        )

    deck = slides_payload.get("pitch_deck") or {}
    raw_slides = deck.get("slides") or []

    if template_path is not None:
        if not template_path.exists():
            raise RuntimeError(f"Template file not found: {template_path}")
        suffix = template_path.suffix.lower()
        if suffix not in {".pptx", ".potx"}:
            raise RuntimeError(
                f"Template must be .pptx or .potx; got '{suffix}' ({template_path})"
            )
        prs = Presentation(str(template_path))
        _remove_existing_slides(prs)
    else:
        prs = Presentation()

    bullets_max = SLIDE_CAPS["bullets_per_slide_max"]
    bullet_chars_max = SLIDE_CAPS["bullet_chars_max"]
    title_chars_max = SLIDE_CAPS["title_chars_max"]
    subtitle_chars_max = SLIDE_CAPS["subtitle_chars_max"]
    notes_chars_max = SLIDE_CAPS["speaker_notes_chars_max"]

    for slide_spec in raw_slides:
        if not isinstance(slide_spec, dict):
            continue
        role = str(slide_spec.get("role") or "")
        title = str(slide_spec.get("title") or "")
        speaker_notes = str(slide_spec.get("speaker_notes") or "")
        layout_name = slide_spec.get("layout_name")

        if role == "title":
            layout = _resolve_layout(prs, layout_name, fallback_index=0)
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = _truncate(title, title_chars_max)
            subtitle = str(slide_spec.get("subtitle") or "")
            if len(slide.placeholders) > 1:
                _set_text_frame(
                    slide.placeholders[1].text_frame,
                    text=subtitle,
                    max_chars=subtitle_chars_max,
                )
            _attach_speaker_notes(slide, speaker_notes, notes_chars_max)
            continue

        bullets_spec = slide_spec.get("bullets") or []
        bullet_texts = [
            str(b.get("text") or "")
            for b in bullets_spec
            if isinstance(b, dict)
        ]
        chunks = _spill_bullets(
            bullets=[{"text": t} for t in bullet_texts],
            bullets_per_slide_max=bullets_max,
        )
        chunk_count = len(chunks)
        for chunk_idx, chunk in enumerate(chunks, start=1):
            layout = _resolve_layout(prs, layout_name, fallback_index=1)
            slide = prs.slides.add_slide(layout)
            slide_title = title
            if chunk_count > 1:
                slide_title = f"{title} ({chunk_idx}/{chunk_count})"
            if slide.shapes.title is not None:
                slide.shapes.title.text = _truncate(slide_title, title_chars_max)
            content_placeholder = _find_content_placeholder(slide)
            if content_placeholder is not None:
                _set_bulleted_text_frame(
                    content_placeholder.text_frame,
                    bullets=[c["text"] for c in chunk],
                    bullet_chars_max=bullet_chars_max,
                )
            _attach_speaker_notes(slide, speaker_notes, notes_chars_max)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _find_content_placeholder(slide):
    """Return the slide's primary content placeholder, or None.

    Title-and-content layouts expose the body placeholder at index 1; some
    template layouts shuffle the index. Iterate to find the first non-title
    placeholder with a text_frame.
    """
    for placeholder in slide.placeholders:
        idx = placeholder.placeholder_format.idx
        if idx == 0:  # title placeholder
            continue
        if placeholder.has_text_frame:
            return placeholder
    return None


def _attach_speaker_notes(slide, notes: str, max_chars: int) -> None:
    if not notes.strip():
        return
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = _truncate(notes, max_chars)


# ---------------------------------------------------------------------------
# Markdown summary — kept around as a sibling artifact to pitch_v{N}.pptx.
# Generated from slides.yaml (so it stays in sync) rather than re-derived
# from artifacts.
# ---------------------------------------------------------------------------


def render_pitch_markdown(slides_payload: dict[str, Any]) -> str:
    deck = slides_payload.get("pitch_deck") or {}
    slides = deck.get("slides") or []
    lines = ["# Stage 4 Pitch", ""]
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        title = slide.get("title") or "(untitled)"
        lines.append(f"## {title}")
        lines.append("")
        if slide.get("role") == "title":
            subtitle = slide.get("subtitle")
            if subtitle:
                lines.append(f"_{subtitle}_")
                lines.append("")
            continue
        for bullet in slide.get("bullets") or []:
            if not isinstance(bullet, dict):
                continue
            text = bullet.get("text") or ""
            lines.append(f"- {text}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"
